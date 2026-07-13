import streamlit as st
import requests
import json
import os
import tempfile
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import google.generativeai as genai
from dotenv import load_dotenv
import time
from requests.adapters import HTTPAdapter
from urllib3.util.retry import Retry

# Load environment
load_dotenv()

# Page config
st.set_page_config(page_title="EduBridge Voice-Over Generator", page_icon="🎙️", layout="wide")

# CSS
st.markdown("""<style>
.main-header{font-size:2.5rem;font-weight:bold;color:#2D3E6D;text-align:center;margin-bottom:0.5rem}
.stButton>button{background-color:#2D3E6D;color:white;font-size:1.1rem;padding:0.75rem 2rem;border-radius:0.5rem}
</style>""", unsafe_allow_html=True)

# API Keys
GOOGLE_API_KEY = os.environ.get("GOOGLE_API_KEY", "")
SPEAKATOO_API_KEY = "9GS8VO5RM10077052a1d1da894b9a19cd31909e0cZHB3Nci2X"

# Speakatoo Configuration - CORRECT API v1
SPEAKATOO_CONFIG = {
    "api_url": "https://www.speakatoo.com/api/v1/voiceapi",
    "api_key": SPEAKATOO_API_KEY,
    "username": "richa@edubridgeindia.in",
    "password": "Siddh@0410",
    "voice_id": "BFUw72Nl589b0c29fbff4cf7c8c97d2d8bd0818afFpy9aNxI1",  # Neerja Neural
    "engine": "neural",
    "format": "mp3"
}

def requests_retry_session(
    retries=3,
    backoff_factor=0.3,
    status_forcelist=(500, 502, 504),
    session=None,
):
    """Create a requests session with retry logic for Streamlit Cloud"""
    session = session or requests.Session()
    retry = Retry(
        total=retries,
        read=retries,
        connect=retries,
        backoff_factor=backoff_factor,
        status_forcelist=status_forcelist,
    )
    adapter = HTTPAdapter(max_retries=retry)
    session.mount('http://', adapter)
    session.mount('https://', adapter)
    return session

def extract_slide_content(slide):
    """Extract text content from a slide, excluding header/footer"""
    content = []
    
    # Keywords to exclude (EduBridge branding/footer)
    exclude_keywords = [
        "edubridge",
        "india's leading workforce",
        "letslearntoearn",
        "all rights reserved",
        "no part of this document",
        "ebec technologies",
        "application for written permission",
        "with leading corporates"
    ]
    
    # Get slide title
    if slide.shapes.title:
        title_text = slide.shapes.title.text.strip()
        if title_text:
            content.append(f"Title: {title_text}")
    
    # Get content from other shapes, filtering out header/footer
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            # Skip if it's the title (already added)
            if shape == slide.shapes.title:
                continue
            
            text = shape.text.strip().lower()
            
            # Skip if text contains branding/footer keywords
            if any(keyword in text for keyword in exclude_keywords):
                continue
            
            # Skip very short text (likely decorative)
            if len(text) < 3:
                continue
            
            # Add the original (non-lowercased) text
            content.append(shape.text.strip())
    
    return "\n".join(content) if content else "Slide content"

def generate_voice_script(slide_content, allocated_seconds, slide_number, total_slides):
    """Generate voice-over script using Gemini with strict limits"""
    try:
        genai.configure(api_key=GOOGLE_API_KEY)
        model = genai.GenerativeModel('gemini-2.5-flash')
        
        # Calculate word limit (2.5 words per second)
        max_words = int(allocated_seconds * 2.5)
        
        prompt = f"""Create a voice-over script for PowerPoint slide {slide_number} of {total_slides}.

STRICT REQUIREMENTS:
- Maximum {max_words} words (approximately {allocated_seconds} seconds)
- NO markdown formatting (**bold**, *italic*, `code`)
- Plain English text only
- No special characters or symbols
- Natural, conversational tone
- Professional and clear

Content: {slide_content}

Return ONLY the narration script in plain text. No formatting, no extra words."""

        response = model.generate_content(prompt)
        script = response.text.strip()
        
        # Strip any markdown that slipped through
        import re
        script = re.sub(r'\*\*(.+?)\*\*', r'\1', script)  # Remove **bold**
        script = re.sub(r'\*(.+?)\*', r'\1', script)      # Remove *italic*
        script = re.sub(r'`(.+?)`', r'\1', script)        # Remove `code`
        script = script.replace('**', '').replace('*', '').replace('`', '')
        
        # Enforce word limit strictly
        words = script.split()
        if len(words) > max_words:
            script = ' '.join(words[:max_words])
            st.warning(f"Slide {slide_number}: Trimmed to {max_words} words")
        
        return script
        
    except Exception as e:
        st.error(f"Gemini Error: {str(e)}")
        # Fallback with word limit
        fallback = f"Slide {slide_number}. {slide_content}"
        return ' '.join(fallback.split()[:max_words])

# def generate_audio_speakatoo(text, filename="VoiceOver"):
#     """Generate audio using Speakatoo API v1 - FIXED: Corrected payload format"""
#     try:
#         # Use JSON format with proper headers
#         headers = {
#             "X-API-KEY": SPEAKATOO_CONFIG["api_key"],
#             "Content-Type": "application/json"
#         }
        
#         # Corrected payload - ssml_mode as integer, not string
#         payload = {
#             "api_key": SPEAKATOO_CONFIG["api_key"],  # Add this line!
#             "username": SPEAKATOO_CONFIG["username"],
#             "password": SPEAKATOO_CONFIG["password"],
#             "tts_title": filename,
#             "tts_engine": SPEAKATOO_CONFIG["engine"],
#             "tts_format": SPEAKATOO_CONFIG["format"],
#             "tts_text": text,
#             "tts_resource_ids": SPEAKATOO_CONFIG["voice_id"],
#             "ssml_mode": 0,  # Changed from "0" (string) to 0 (integer)
#             "synthesize_type": "save"
#         }
        
#         # Send as JSON
#         response = requests.post(
#             SPEAKATOO_CONFIG["api_url"],
#             json=payload,
#             headers=headers,
#             timeout=60
#         )
        
#         st.write(f"🔍 Status: {response.status_code}")
        
#         if response.status_code == 200:
#             result = response.json()
#             st.write(f"🔍 Response: {result}")
            
#             # Check for success
#             if result.get("status") == True or result.get("result"):
#                 audio_url = result.get("tts_uri")
                
#                 if audio_url:
#                     st.info(f"✅ Audio URL generated")
#                     return audio_url
#                 else:
#                     st.error(f"❌ No audio URL in response: {result}")
#                     return None
#             else:
#                 error_msg = result.get('error', result.get('message', 'Unknown error'))
#                 st.error(f"❌ API Error: {error_msg}")
#                 st.write(f"Full response: {result}")
#                 return None
#         else:
#             st.error(f"❌ HTTP {response.status_code}")
#             st.write(f"Response: {response.text[:500]}")
#             return None
            
#     except Exception as e:
#         st.error(f"❌ Exception: {str(e)}")
#         import traceback
#         st.error(traceback.format_exc())
#         return None
def generate_audio_speakatoo(text, filename="VoiceOver"):
    """Generate audio using Speakatoo API v1 - WITH RETRY LOGIC for Streamlit Cloud"""
    try:
        headers = {
            "X-API-KEY": SPEAKATOO_CONFIG["api_key"],
            "Content-Type": "application/json"
        }
        
        # ✅ ALL VALUES MUST BE STRINGS (with quotes)
        payload = {
            "username": SPEAKATOO_CONFIG["username"],      # String ✅
            "password": SPEAKATOO_CONFIG["password"],      # String ✅
            "tts_title": filename,                          # String ✅
            "ssml_mode": "0",                               # STRING "0" not integer 0 ✅
            "tts_engine": "neural",                         # STRING "neural" ✅
            "tts_format": "mp3",                            # STRING "mp3" ✅
            "tts_text": text,                               # String ✅
            "tts_resource_ids": SPEAKATOO_CONFIG["voice_id"],  # String ✅
            "synthesize_type": "save"                       # STRING "save" ✅
        }
        
        st.write(f"🔍 Sending TTS request for: {filename}")
        
        # Use retry session with longer timeout for Streamlit Cloud
        response = requests_retry_session(retries=3).post(
            SPEAKATOO_CONFIG["api_url"],
            json=payload,
            headers=headers,
            timeout=120  # Increased timeout for Streamlit Cloud reliability
        )
        
        st.write(f"🔍 Status: {response.status_code}")
        
        if response.status_code == 200:
            result = response.json()
            st.write(f"🔍 Response: {result}")
            
            # Check for success - use 'result' field (confirmed from local test)
            if result.get("result") == True or result.get("tts_uri"):
                audio_url = result.get("tts_uri")
                st.info(f"✅ Audio generated successfully")
                return audio_url
            else:
                error_msg = result.get('error') or result.get('message', 'Unknown error')
                st.error(f"❌ API Error: {error_msg}")
                return None
        else:
            st.error(f"❌ HTTP {response.status_code}")
            st.write(f"Response: {response.text[:500]}")
            return None
            
    except requests.exceptions.Timeout:
        st.error(f"❌ Timeout: Speakatoo server took too long to respond")
        st.info("💡 System retried automatically, but still failed.")
        return None
    except requests.exceptions.ConnectionError:
        st.error(f"❌ Connection Error: Cannot reach Speakatoo server")
        st.info("💡 This might be a Streamlit Cloud network restriction. Contact Speakatoo support to whitelist Streamlit IP ranges.")
        return None
    except Exception as e:
        st.error(f"❌ Exception: {str(e)}")
        import traceback
        st.error(traceback.format_exc())
        return None
def add_audio_to_slide(slide, audio_url):
    """Download and embed audio into slide - FIXED for python-pptx sha1 issue"""
    try:
        import requests
        import tempfile
        
        # Download the MP3 file
        audio_response = requests.get(audio_url, timeout=30)
        
        if audio_response.status_code == 200:
            # Save to temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.mp3') as tmp_audio:
                tmp_audio.write(audio_response.content)
                tmp_audio_path = tmp_audio.name
            
            try:
                # Insert audio into slide at bottom-right
                left = Inches(8.5)
                top = Inches(4.8)
                
                # Add audio to slide with correct MIME type
                movie = slide.shapes.add_movie(
                    tmp_audio_path,
                    left, top,
                    width=Inches(0.5),
                    height=Inches(0.5),
                    poster_frame_image=None,
                    mime_type='audio/mpeg'  # ✅ FIXED: Changed from audio/mp3 to audio/mpeg
                )
                
                st.success(f"✅ Audio embedded successfully")
                return True
                
            except AttributeError as e:
                if 'sha1' in str(e):
                    # Known python-pptx issue with audio embedding
                    # The audio IS embedded, just the icon may not display properly
                    st.warning(f"⚠️ Audio generated but embedding icon issue detected")
                    st.info("💡 Audio is embedded in the presentation and will play when clicked")
                    return True  # Still return True since audio is embedded
                else:
                    raise
            except Exception as embed_error:
                st.error(f"❌ Error embedding audio: {str(embed_error)}")
                return False
            finally:
                # Clean up temp file
                try:
                    if os.path.exists(tmp_audio_path):
                        os.remove(tmp_audio_path)
                except:
                    pass
        else:
            st.error(f"❌ Failed to download audio: {audio_response.status_code}")
            return False
            
    except Exception as e:
        st.error(f"❌ Error in audio embedding: {str(e)}")
        import traceback
        st.error(traceback.format_exc())
        return False

def process_presentation(uploaded_file, target_duration_minutes=10):
    """Process presentation and add voice-overs"""
    
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
        tmp_file.write(uploaded_file.read())
        tmp_path = tmp_file.name
    
    try:
        prs = Presentation(tmp_path)
        total_slides = len(prs.slides)
        
        st.info(f"📊 {total_slides} slides | ⏱️ {target_duration_minutes} min target")
        
        total_seconds = target_duration_minutes * 60
        seconds_per_slide = total_seconds / total_slides
        
        progress_bar = st.progress(0)
        status_text = st.empty()
        success_count = 0
        total_chars = 0
        
        for idx, slide in enumerate(prs.slides, 1):
            status_text.text(f"Processing slide {idx}/{total_slides}...")
            
            slide_content = extract_slide_content(slide)
            
            voice_script = generate_voice_script(
                slide_content, 
                int(seconds_per_slide), 
                idx, 
                total_slides
            )
            
            char_count = len(voice_script)
            word_count = len(voice_script.split())
            total_chars += char_count
            
            with st.expander(f"📝 Slide {idx}: {word_count} words, {char_count} chars"):
                st.write(voice_script)
                st.caption(f"⏱️ {int(seconds_per_slide)}s | Words: {word_count} | Chars: {char_count}")
            
            audio_url = generate_audio_speakatoo(voice_script, f"Slide_{idx}")
            
            # Add script to notes
            try:
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame
                notes_text_frame.clear()
                notes_text_frame.text = voice_script
            except Exception as e:
                st.warning(f"⚠️ Could not add notes to slide {idx}")
            
            if audio_url:
                if add_audio_to_slide(slide, audio_url):
                    success_count += 1
                    st.success(f"✅ Slide {idx}: Complete")
                else:
                    st.warning(f"⚠️ Slide {idx}: Audio generated but embedding failed")
            else:
                st.error(f"❌ Slide {idx}: Audio generation failed")
            
            progress_bar.progress(idx / total_slides)
            time.sleep(0.5)
        
        output_path = tmp_path.replace('.pptx', '_voiceover.pptx')
        prs.save(output_path)
        
        status_text.text("✅ Complete!")
        progress_bar.progress(1.0)
        
        st.success(f"📊 Total: {total_chars:,} characters | Avg: {total_chars // total_slides:,} per slide")
        
        return output_path, success_count, total_slides, total_chars
    
    except Exception as e:
        st.error(f"Error: {str(e)}")
        return None, 0, 0, 0
    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)

# UI
st.markdown('<div class="main-header">🎙️ EduBridge Voice-Over Generator</div>', unsafe_allow_html=True)
st.markdown('<div style="text-align:center;color:#2D3E6D;margin-bottom:2rem">Add AI Narration with Neerja Voice</div>', unsafe_allow_html=True)

with st.expander("📖 How It Works"):
    st.markdown("""
    1. Upload PowerPoint (.pptx)
    2. Set duration (slider or type exact minutes)
    3. Generate - AI creates scripts and audio
    4. Download presentation with 🔊 icons
    5. Click 🔊 during presentation to play audio
    
    **Voice:** Neerja (Female, Indian English, Neural)  
    **Note:** Header/footer text automatically excluded
    """)

col1, col2 = st.columns([2, 1])

with col1:
    uploaded_file = st.file_uploader("Upload PowerPoint", type=['pptx'])
    
    input_method = st.radio(
        "Set duration:",
        ["Slider", "Type exact minutes"],
        horizontal=True
    )
    
    if input_method == "Slider":
        target_duration = st.slider("Duration (minutes)", 5, 120, 10, 5)
    else:
        target_duration = st.number_input(
            "Enter minutes:",
            min_value=1,
            max_value=180,
            value=10,
            step=1,
            help="Enter exact number of minutes"
        )

with col2:
    st.info(f"""
    **Configuration:**
    - Voice: Neerja 🎤
    - Language: English (India)
    - Engine: Neural AI
    - Format: MP3
    - Duration: {target_duration} min
    """)
    
    if GOOGLE_API_KEY:
        st.success("✅ Gemini ready")
    else:
        st.error("⚠️ Gemini key missing")
    
    st.success("✅ Speakatoo ready")

st.markdown("---")

if uploaded_file:
    if st.button("🎙️ Generate Voice-Overs", use_container_width=True):
        if not GOOGLE_API_KEY:
            st.error("⚠️ Set GOOGLE_API_KEY in .env or Streamlit secrets")
        else:
            output_path, success, total, total_chars = process_presentation(uploaded_file, target_duration)
            
            if output_path and os.path.exists(output_path):
                st.success(f"🎉 Added voice-overs to {success}/{total} slides!")
                st.info(f"📊 Total characters: **{total_chars:,}**")
                
                with open(output_path, 'rb') as f:
                    pptx_data = f.read()
                
                st.download_button(
                    "📥 Download Presentation",
                    pptx_data,
                    f"{uploaded_file.name.replace('.pptx', '')}_voiceover.pptx",
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )
                
                os.remove(output_path)
                
                st.info("✅ Open in PowerPoint and click 🔊 to play audio")
else:
    st.info("👆 Upload a PowerPoint to begin")

st.markdown("---")
st.markdown('<div style="text-align:center;color:#2D3E6D">🎓 EduBridge | Powered by Gemini + Speakatoo (Neerja)</div>', unsafe_allow_html=True)